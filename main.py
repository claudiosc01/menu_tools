from pystyle import Colorate, Colors, Write, Add, Center
import os
import requests
import time
import random
import string
from fpdf import FPDF
from pypresence import Presence
import random
from pptx import Presentation
from console_progressbar import ProgressBar
from bs4 import BeautifulSoup
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

def Banner():
    os.system('cls')
    Banner1 = r"""
█████  █████ ██████   ██████ ███████████   █████████  █████   █████  ████████   ████████   ████████ 
░░███  ░░███ ░░██████ ██████ ░░███░░░░░███ ███░░░░░███░░███   ░░███  ███░░░░███ ███░░░░███ ███░░░░███
 ░███   ░███  ░███░█████░███  ░███    ░███░███    ░░░  ░███    ░███ ░███   ░░░ ░███   ░░░ ░███   ░░░ 
 ░███   ░███  ░███░░███ ░███  ░██████████ ░░█████████  ░███████████ ░█████████ ░█████████ ░█████████ 
 ░███   ░███  ░███ ░░░  ░███  ░███░░░░░░   ░░░░░░░░███ ░███░░░░░███ ░███░░░░███░███░░░░███░███░░░░███
 ░███   ░███  ░███      ░███  ░███         ███    ░███ ░███    ░███ ░███   ░███░███   ░███░███   ░███
 ░░████████   █████     █████ █████       ░░█████████  █████   █████░░████████ ░░████████ ░░████████ 
  ░░░░░░░░   ░░░░░     ░░░░░ ░░░░░         ░░░░░░░░░  ░░░░░   ░░░░░  ░░░░░░░░   ░░░░░░░░   ░░░░░░░░  
                                                                                                                 
        @umplsh666
"""

    Banner2 = r"""
  ,           ,
 /             \
((__-^^-,-^^-__))
 `-_---' `---_-'
  <__|o` 'o|__>
     \  `  /
      ): :(
      :o_o:
       "-" 
       """

    print(Center.XCenter(Colorate.Vertical(Colors.yellow_to_red, Add.Add(Banner2, Banner1, center=True), 2)))


def Banner_end():
    os.system('cls')
    Banner1 = u"\n\n                       ♛ Finally. Good H4cking! ;)"

    Banner2 = r"""
  ,           ,
 /  @umpsh666  \
((__-^^-,-^^-__))
 `-_---' `---_-'
  <__|o` 'o|__>
     \  `  /
      ): :(
      :o_o:
       "-" 
       """

    print(Center.XCenter(Colorate.Vertical(Colors.yellow_to_red, Add.Add(Banner1, Banner2, center=True), 2)))

    
#-----------------------------------------------------------------------------------------------------
def values_texts_options():
    config_values1 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 1. Crear proyecto MVC (NODE.JS) ")
    config_values2 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 2. Obtener estructura proyecto. ")
    config_values3 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 3. Crear proyecto MVC (EXPRESS.JS) API. ")
    config_values4 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 4. Generar PPTXS random para subir a paginas. ")
    config_values5 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 5. Generar PDFS random para subir a paginas. ")
    config_values6 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 6. Generador de Discord Nitro.                           [Slower]")
    config_values7 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 7. Obtener informacion de WhatsApp.                      [In Process]")
    config_values8 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 8. Gif-Imagen para WhatsApp v1. ")
    config_values9 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 9. Generar Views para Youtube.                           [In Process]")
    config_values10 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 10. Enviar multimensajes a tus Contactos - WhatsApp v1.  [Inactive]")
    config_values11 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 11. Gif-Imagen para Discord v1.  [Slower]")
    config_values12 = Colorate.Horizontal(Colors.blue_to_cyan, "      > 12. Personalización de presencia de Discord - (Discord RPC)")
    print('\n\n'+config_values1+'\n'+config_values2+'\n'+config_values3+'\n'+config_values4+'\n'+config_values5+'\n'+config_values6+'\n'+config_values7+'\n'+config_values8+'\n'+config_values9+'\n'+config_values10+'\n'+config_values11+'\n'+config_values12)

#---------------------------------------------------FUNCTIONS--------------------------------------------------------




#-----------------------------------------------------------------------------------------------------

def generar_estructura_proyecto_api(directorio_base, nombre_proyecto):
    # Generar la ruta de la carpeta raíz del proyecto
    root = os.path.join(directorio_base, nombre_proyecto)
    
    # Definir las subcarpetas
    dirs = [
        'node_modules',
        'public',
        'src/controllers',
        'src/middlewares',
        'src/models',
        'src/routes',
        'src/services',
        'test',
    ]

    # Crear cada subcarpeta dentro de la carpeta raíz del proyecto
    for dir in dirs:
        os.makedirs(os.path.join(root, dir), exist_ok=True)

    # Definir los archivos a crear
    files = [
        '.env',
        '.gitignore',
        'src/index.js',
    ]

    # Crear cada archivo dentro de la carpeta raíz del proyecto
    for file in files:
        open(os.path.join(root, file), 'a').close()

    os.chdir(root)
    os.system('npm init -y')
    os.system('cls')

#-----------------------------------------------------------------------------------------------------





#-----------------------------------------------------------------------------------------------------
def generar_estructura_proyecto(directorio_base, nombre_proyecto):
    
    # Generar la ruta de la carpeta raíz del proyecto
    root = os.path.join(directorio_base, nombre_proyecto)
    
    # Definir las subcarpetas
    dirs = [
        'config',
        'controllers',
        'helpers',
        'models',
        'routes',
        'views',
    ]

    # Crear cada subcarpeta dentro de la carpeta raíz del proyecto
    for dir in dirs:
        os.makedirs(os.path.join(root, dir), exist_ok=True)

    # Definir los archivos a crear
    files = [
        '.env',
        'index.js',
        'config/db.js',
    ]

    # Crear cada archivo dentro de la carpeta raíz del proyecto
    for file in files:
        open(os.path.join(root, file), 'a').close()

#-----------------------------------------------------------------------------------------------------





#-----------------------------------------------------------------------------------------------------
def listar_contenido(directorio):
    for root, dirs, files in os.walk(directorio):
        # Ignorar el directorio node_modules/
        if 'node_modules' in dirs:
            dirs.remove('node_modules')
            
        level = root.replace(directorio, '').count(os.sep)
        indent = ' ' * 4 * level
        print('----{}{}/'.format(indent, os.path.basename(root))) #get sub-carpetas
        subindent = ' ' * 4 * (level + 1)
        for f in files:
            print('++++{}{}'.format(subindent, f)) #get -archivos
#-----------------------------------------------------------------------------------------------------
            

def generate_random_presentation(directorio):
    if not os.path.exists(directorio):
        os.makedirs(directorio)

    for i in range(7):
        prs = Presentation()

        for j in range(50):
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            texto_random = ''.join(random.choices(string.ascii_letters + string.digits, k=20))
            title.text = f"Página {j+1}"
            subtitle.text = texto_random

        archivo_pptx = os.path.join(directorio, f"documento0{i+1}.pptx")
        prs.save(archivo_pptx)


    response_pptxs = Colorate.Horizontal(Colors.cyan_to_green, "[+] Se han creado correctamente los PPTXS. Good D4y :3 ")
    print(f"\n{response_pptxs}")

    respuesta_input = Colorate.Horizontal(Colors.cyan_to_green, "                    | ¿Desea eliminar todos los archivos en este directorio? (s/n) > ")
    respuesta = str(input(respuesta_input))


    if respuesta.lower() in {'s', 'si'}:
        wait_delete = Colorate.Horizontal(Colors.yellow_to_red, "    [+] Espera a la eliminacion de todos los archivos contenidos del directorio.")
        print(f"\n{wait_delete}")
        time.sleep(3)
        for filename in os.listdir(directorio):
            if filename.endswith(".pptx"):
                os.remove(os.path.join(directorio, filename))
        correct_delete = Colorate.Horizontal(Colors.yellow_to_red, "    [+] Todos los archivos de presentación en el directorio han sido eliminados. Good D4y :3")
        print(f"{correct_delete}")
        time.sleep(2)
    else:
        error_delete = Colorate.Horizontal(Colors.yellow_to_red, "    [!] Los archivos no han sido eliminados. Intentalo Manualmente. Good D4y :3")
        print(f"\n{error_delete}")
        time.sleep(2)




#-----------------------------------------------------------------------------------------------------


def generate_random_pdf(directorio):
    if not os.path.exists(directorio):
        os.makedirs(directorio)

    for i in range(7):
        pdf = FPDF()

        # Ajustes de la página
        pdf.add_page()

        # Configuración de la fuente
        pdf.set_font("Arial", size=15)

        # Generar 50 líneas de texto aleatorio
        for j in range(100):
            texto_random = ''.join(random.choices(string.ascii_letters + string.digits, k=20))
            pdf.cell(200, 10, txt=texto_random, ln=1, align='C')

        archivo_pdf = os.path.join(directorio, f"documento0{i + 1}.pdf")
        pdf.output(archivo_pdf)

    response_pdfs = Colorate.Horizontal(Colors.cyan_to_green, "[+] Se han creado correctamente los PDFs. Good D4y :3 ")
    print(f"\n{response_pdfs}")

    respuesta_input = Colorate.Horizontal(Colors.cyan_to_green, "                    | ¿Desea eliminar todos los archivos en este directorio? (s/n) > ")
    respuesta = str(input(respuesta_input))

    if respuesta.lower() in {'s', 'si'}:
        wait_delete = Colorate.Horizontal(Colors.yellow_to_red, "    [+] Espera a la eliminacion de todos los archivos contenidos del directorio.")
        print(f"\n{wait_delete}")
        time.sleep(3)
        for filename in os.listdir(directorio):
            if filename.endswith(".pdf"):
                os.remove(os.path.join(directorio, filename))
        correct_delete = Colorate.Horizontal(Colors.yellow_to_red, "    [+] Todos los archivos PDF en el directorio han sido eliminados. Good D4y :3")
        print(f"{correct_delete}")
        time.sleep(2)
    else:
        error_delete = Colorate.Horizontal(Colors.yellow_to_red, "    [!] Los archivos no han sido eliminados. Intentalo Manualmente. Good D4y :3")
        print(f"\n{error_delete}")
        time.sleep(2)


#-----------------------------------------------------------------------------------------------------


chrome_options = Options()
chrome_options.add_argument("--log-level=3")

def discord_checker_nitro(number_range):
    driver = webdriver.Chrome(options=chrome_options)
    message_wait = Colorate.Horizontal(Colors.cyan_to_green, "              | Generando discord nitro... :,)")
    print(f"{message_wait}\n")

    for x in range(number_range):
        code = "".join(random.choice("abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ1234567890") for _ in range(16))
        
        driver.get(f'https://discord.com/gifts/{code}')
        time.sleep(6)
        element = WebDriverWait(driver, 10).until(
            EC.presence_of_element_located((By.CSS_SELECTOR, "div.centeringWrapper__319b0"))
        )

        html = driver.page_source
        soup = BeautifulSoup(html, 'html.parser')

        h1_text = soup.find('div', class_='centeringWrapper__319b0').text
        if h1_text == "Código de regalo no válidoPuede que este código de regalo haya caducado o quizás tengas un código incorrecto.Continuar en Discord¿Por qué mi código de regalo no es válido?":
            data_boost_false_discord = Colorate.Horizontal(Colors.yellow_to_red, f"     > Discord Link https://discord.com/gifts/{code} - False Boost Discord")
            print(data_boost_false_discord)
        else:
            data_boost_true_discord = Colorate.Horizontal(Colors.cyan_to_green, f"     > Discord Link https://discord.com/gifts/{code} - True Boost Discord")
            print(data_boost_true_discord)
        
        time.sleep(3)

    close_func_discord = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n > Press enter to closed...")
    input(close_func_discord)
    exit

#-----------------------------------------------------------------------------------------------------
    
def whatsapp_getting_chats():
    driver = webdriver.Chrome(options=chrome_options)
    driver.get('https://web.whatsapp.com/')

    input("Una vez hayas escaneado el código QR y hecho inicio de sesión, presiona Enter.")
    driver.get('https://web.whatsapp.com/')
    time.sleep(11)
    # print(driver.page_source)


    element = WebDriverWait(driver, 10).until(
        EC.presence_of_element_located((By.CSS_SELECTOR, "div.lhggkp7q"))
    )

    html = driver.page_source
    soup = BeautifulSoup(html, 'html.parser')
    h1_text = soup.find('div', class_='_3YS_f').text
    print(h1_text)

    #Falta obtener los chats, numeros, 


    # enter_to_run_getter = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n  > Press enter if you logged in correctly WhatsApp-Web...")
    # input(enter_to_run_getter)
    close_func_discord = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n > Press enter to closed...")
    input(close_func_discord)
    exit

#-----------------------------------------------------------------------------------------------------


def whatsapp_gifter_change_photo():
    driver = webdriver.Chrome(options=chrome_options)
    driver.get('https://web.whatsapp.com/')

    enter_to_run_getter = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n  > Press enter if you logged in correctly WhatsApp-Web...")
    input(enter_to_run_getter)
    
    driver.get('https://web.whatsapp.com/')
    time.sleep(11)
    # print(driver.page_source)

    driver.find_element(By.CSS_SELECTOR, 'div.g0rxnol2').click() 
    time.sleep(2)

    message_wsp_closed = Colorate.Horizontal(Colors.yellow_to_red, f"  > If u need exit to program. Prees Ctrl + C")
    print(message_wsp_closed)
    try:
        while True:
            for x in range(5):
                # print(f"> range {x}")
                if x % 2 == 0:
                    upload = driver.find_element(By.XPATH, '//input[@type="file"]')
                    upload.send_keys('C:/Users/parak/Pictures/sad.jpg')
                    time.sleep(0.4)
                    driver.find_element(By.CSS_SELECTOR, 'div._3oDXB').click()
                else:
                    upload = driver.find_element(By.XPATH, '//input[@type="file"]')
                    upload.send_keys('C:/Users/parak/Pictures/happy.jpg')
                    time.sleep(0.4)
                    driver.find_element(By.CSS_SELECTOR, 'div._3oDXB').click()
            driver.get('https://web.whatsapp.com/')
            time.sleep(60)
            driver.find_element(By.CSS_SELECTOR, 'div.g0rxnol2').click() 
            time.sleep(2)
    except:
        print("An exception occurred")
        close_func_discord = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n > Press enter to closed...")
        input(close_func_discord)
        exit

#-----------------------------------------------------------------------------------------------------

def discord_gifter_change_photo():
    driver = webdriver.Chrome(options=chrome_options)
    driver.get('https://discord.com/channels/@me')

    enter_to_run_getter = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n  > Press enter if you logged in correctly Discord-Web...")
    input(enter_to_run_getter)
    driver.get('https://discord.com/channels/@me')
    

    time.sleep(4)
    elements = driver.find_elements(By.CLASS_NAME, 'contents_fb6220') #Pagina Principal, boton de opcions - 1

    if len(elements) > 2:
        elements[2].click()
        time.sleep(2)
    
    driver.find_element(By.CLASS_NAME, 'button_afdfd9.lookFilled__19298.colorBrand_b2253e.sizeSmall__71a98.grow__4c8a4').click()
    time.sleep(3)

    try:
        while True:
            for x in range(3):
                if x % 2 == 0:
                    driver.find_element(By.CLASS_NAME, 'wrapper_edb6e0.avatarUploaderInner_c81617').click()

                    #cambiar avatar
                    wait = WebDriverWait(driver, 2)  
                    elements2 = wait.until(EC.visibility_of_all_elements_located((By.CLASS_NAME, 'label__73cb9')))
                    if elements2:
                        elements2[0].click()

                    # file-input
                    wait = WebDriverWait(driver, 10)  # Espera hasta 10 segundos
                    upload = wait.until(EC.presence_of_element_located((By.XPATH, '//*[@id="app-mount"]/div[2]/div[1]/div[4]/div[2]/div/div/div[2]/div[1]/div[1]/input')))
                    upload.send_keys('C:/Users/parak/Pictures/arcoris-perfil.png')
                    
                    time.sleep(3)
                    # Accept Photo Perfil
                    wait = WebDriverWait(driver, 5)  # Esperar hasta 10 segundos
                    elements3 = wait.until(EC.presence_of_all_elements_located((By.CSS_SELECTOR, '.button_afdfd9.lookFilled__19298.colorBrand_b2253e.sizeSmall__71a98.grow__4c8a4')))

                    if len(elements3) > 7:  # Asegurarse de que hay al menos 8 elementos
                        elements3[7].click()
                    time.sleep(2)
                    wait = WebDriverWait(driver, 2)
                    elements4 = wait.until(EC.visibility_of_all_elements_located((By.CSS_SELECTOR, '.button_afdfd9.lookFilled__19298.colorGreen__5f181.sizeSmall__71a98.grow__4c8a4')))
                    if elements4: 
                        elements4[0].click()  
                    
                    time.sleep(30)

                else:
                    driver.find_element(By.CLASS_NAME, 'wrapper_edb6e0.avatarUploaderInner_c81617').click()

                    #cambiar avatar
                    wait = WebDriverWait(driver, 2)
                    elements2 = wait.until(EC.visibility_of_all_elements_located((By.CLASS_NAME, 'label__73cb9')))
                    if elements2:
                        elements2[0].click()

                    # file-input
                    wait = WebDriverWait(driver, 5) 
                    upload = wait.until(EC.presence_of_element_located((By.XPATH, '//*[@id="app-mount"]/div[2]/div[1]/div[4]/div[2]/div/div/div[2]/div[1]/div[1]/input')))
                    upload.send_keys('C:/Users/parak/Pictures/lel.png')
                    
                    time.sleep(3)
                    # Accept Photo Perfil
                    wait = WebDriverWait(driver, 5)  # Esperar hasta 10 segundos
                    elements3 = wait.until(EC.presence_of_all_elements_located((By.CSS_SELECTOR, '.button_afdfd9.lookFilled__19298.colorBrand_b2253e.sizeSmall__71a98.grow__4c8a4')))

                    if len(elements3) > 7:  # Asegurarse de que hay al menos 8 elementos
                        elements3[7].click()  

                    time.sleep(2)
                    wait = WebDriverWait(driver, 2)
                    elements4 = wait.until(EC.visibility_of_all_elements_located((By.CSS_SELECTOR, '.button_afdfd9.lookFilled__19298.colorGreen__5f181.sizeSmall__71a98.grow__4c8a4')))
                    if elements4:  
                        elements4[0].click()  


                    time.sleep(30)
                    
            time.sleep(5)
            driver.get('https://discord.com/channels/@me')
            elements = driver.find_elements(By.CLASS_NAME, 'contents_fb6220') #Pagina Principal, boton de opcions - 1

            if len(elements) > 2:
                elements[2].click()
                time.sleep(2)
            
            driver.find_element(By.CLASS_NAME, 'button_afdfd9.lookFilled__19298.colorBrand_b2253e.sizeSmall__71a98.grow__4c8a4').click()
            time.sleep(60)

    except:
        print("An exception occurred")
        close_func_discord = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n > Press enter to closed...")
        input(close_func_discord)
        exit

#-----------------------------------------------------------------------------------------------------
def generate_views_for_youtube(views_number, link_video, seconds):
    chrome_options.add_argument("--headless")
    chrome_options.add_argument("--mute-audio")
    driver = webdriver.Chrome(options=chrome_options)
    count_views=0

        
    while count_views != views_number:
        print(count_views)

        # Cierra todas las pestañas menos la primera
        for handle in driver.window_handles[1:]:
            driver.switch_to.window(handle)
            driver.close()

        # Cambia a la primera pestaña
        driver.switch_to.window(driver.window_handles[0])

        for x in range(3):
            driver.execute_script(f"window.open('{link_video}');")
            driver.switch_to.window(driver.window_handles[-1])
            if x%3 == 1:    
                time.sleep(6)
                driver.find_element(By.CSS_SELECTOR, '.ytp-mute-button').click()
            else:
                time.sleep(2)
                pass
        time.sleep(seconds)
        driver.close()
        count_views+=3

    message_succelly = Colorate.Horizontal(Colors.cyan_to_green, f"\n\n > Se han generado correctamente las {count_views} vistas. :,)")
    print(message_succelly)
    
    close_func_discord = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n > Press enter to closed...")
    input(close_func_discord)
    exit
#-----------------------------------------------------------------------------------------------------
    
def generate_automessages_whatsapp(number):
    driver = webdriver.Chrome(options=chrome_options)
    driver.get('https://web.whatsapp.com/')

    enter_to_run_getter = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n  > Press enter if you logged in correctly WhatsApp-Web...")
    input(enter_to_run_getter)
    
    driver.get('https://web.whatsapp.com/')
    time.sleep(11)
    # print(driver.page_source)
    # 1. Clicked in the search.
    # 2. Paste the number in the search.
    # 3. Clicked in the message tabular or press enter for send message.
    # 4. Paste the message for the contact.
    # 5. In Bucle for other numbers of the WhatsApp.
    pass

#-----------------------------------------------------------------------------------------------------
def get_current_time(start_time):
    # Calcula el tiempo transcurrido en segundos desde el inicio
    elapsed_seconds = time.time() - start_time
    # Calcula las horas, minutos y segundos restantes
    remaining_seconds = 99 * 3600 + elapsed_seconds
    hours = int(remaining_seconds // 3600)
    minutes = int((remaining_seconds % 3600) // 60)
    seconds = int(remaining_seconds % 60)
    # Formatea el tiempo en formato HH:MM:SS como una cadena de texto
    return "{:02d}:{:02d}:{:02d}".format(hours, minutes, seconds)

# Tiempo actual... 24hours
start_time = time.time()


def discord_rpc_open(id_client,urlImage1,urlImage2,description,buttonsTexts=[]):

    try:

        # Crear una instancia de Presence
        RPC = Presence(id_client)

        # Conectar al cliente de Discord
        RPC.connect()

        # Bucle infinito para actualizar continuamente la presencia
        while True:
            # Actualizar la presencia con la hora actual en el estado y la imagen de fondo
            RPC.update(
                state="Activity elapsed: " + get_current_time(start_time),
                large_image=urlImage1,
                small_image=urlImage2,
                details="Activity Summary:   " + description, 
                # end=1+1,
                start=1,
                # end=10,
                # pid=114297,
                # party_id="1jdhg5ks",
                party_size=[1,1],
                buttons=buttonsTexts
                # details=get_current_time(),
            )

            # Esperar un intervalo de tiempo (en segundos)
            time.sleep(1)  # Por ejemplo, actualizar cada 1 segundos

    except KeyboardInterrupt:
        # Cerrar la conexión cuando se presione Ctrl+C
        RPC.close()
#-----------------------------------------------------------------------------------------------------


#-----------------------------------------------------------------------------------------------------
def responses_functions():
    values_texts_options()
    config_text = Colorate.Horizontal(Colors.yellow_to_red, f"\n\n | Ingresa la opcion que quieres realizar >  ")
    config_data = int(input(config_text))

    if config_data == 1:
        directorio_base_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingrese el directorio base > ")
        directorio_base = str(input(directorio_base_prompt))

        nombre_proyecto_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                                           | Ingrese el nombre del proyecto > ")
        nombre_proyecto = str(input(nombre_proyecto_prompt))

        generar_estructura_proyecto(directorio_base, nombre_proyecto)
    elif config_data == 2:
        setDirectory_prompt = Colorate.Horizontal(Colors.green_to_blue, "               | Ingrese el path del directorio > ")
        setDirectory = str(input(setDirectory_prompt))
        directorio_raiz = setDirectory
        listar_contenido(directorio_raiz)
    elif config_data == 3:
        directorio_base_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingrese el directorio base > ")
        directorio_base = str(input(directorio_base_prompt))
        nombre_proyecto_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                                           | Ingrese el nombre del proyecto-api > ")
        nombre_proyecto = str(input(nombre_proyecto_prompt))
        generar_estructura_proyecto_api(directorio_base, nombre_proyecto)
    elif config_data == 4:
        directorio_base_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingrese el directorio base > ")
        directorio_base = str(input(directorio_base_prompt))
        generate_random_presentation(directorio_base)
    elif config_data == 5:
        directorio_base_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingrese el directorio base > ")
        directorio_base = str(input(directorio_base_prompt))
        generate_random_pdf(directorio_base)
    elif config_data == 6:
        number_generate_discord_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa el numero de discord a generar: > ")
        number_generate_discord = int(input(number_generate_discord_prompt))
        discord_checker_nitro(number_generate_discord)
    elif config_data == 7:
        whatsapp_getting_chats()
    elif config_data == 8:
        whatsapp_gifter_change_photo()
    elif config_data == 9:
        paste_link_for_generate_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa el link del video: > ")
        paste_link_for_generate = str(input(paste_link_for_generate_prompt))

        generate_views_number_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa el numero de vistas a generar: > ")
        generate_views_number = int(input(generate_views_number_prompt))

        generate_seconds_for_link_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa los segundos para cada video: > ")
        generate_seconds_for_link = int(input(generate_seconds_for_link_prompt))

        generate_views_for_youtube(generate_views_number, paste_link_for_generate, generate_seconds_for_link)
    elif config_data == 10:
        search_contact_number_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa el numero del contacto de WhatsApp: > ")
        search_contact_number = str(input(paste_link_for_generate_prompt))
        generate_automessages_whatsapp()
    elif config_data == 11:
        discord_gifter_change_photo()
    elif config_data == 12:

        mainTexts_buttons=[]


        idClientDiscord_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa el Id-Client de tu Bot Creado: > ")
        idClientDiscord = str(input(idClientDiscord_prompt))

        imageLarge_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa la primera imagen grande a mostrar: > ")
        imageLarge = str(input(imageLarge_prompt))

        imageSmall_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa la  imagen pequeña a mostrar: > ")
        imageSmall = str(input(imageSmall_prompt))

        
        description_text_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingrese una descripcion pequeña: > ")
        description_text = str(input(description_text_prompt))


        mainTextBtn1_label_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa el texto a mostrar para el boton 1: > ")
        mainTextBtn1_url1_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa la url a redireccionar para el boton 1: > ")
        mainTextBtn1_label = str(input(mainTextBtn1_label_prompt))
        mainTextBtn1_url1 = str(input(mainTextBtn1_url1_prompt))

        buttom_first = {"label": "> " + mainTextBtn1_label, "url": mainTextBtn1_url1}
        mainTexts_buttons.append(buttom_first)

        mainTextBtn2_label_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa el texto a mostrar para el boton 2: > ")
        mainTextBtn2_url1_prompt = Colorate.Horizontal(Colors.cyan_to_green, "                            | Ingresa la url a redireccionar para el boton 2: > ")
        mainTextBtn2_label = str(input(mainTextBtn2_label_prompt))
        mainTextBtn2_url2 = str(input(mainTextBtn2_url1_prompt))
        
        buttom_second = {"label": "> " + mainTextBtn2_label, "url": mainTextBtn2_url2}
        mainTexts_buttons.append(buttom_second)

        discord_rpc_open(idClientDiscord,imageLarge,imageSmall,description_text,mainTexts_buttons)

        configuration = '''import time
from pypresence import Presence   

def get_current_time(start_time):
    # Calcula el tiempo transcurrido en segundos desde el inicio
    elapsed_seconds = time.time() - start_time
    # Calcula las horas, minutos y segundos restantes
    remaining_seconds = 99 * 3600 + elapsed_seconds
    hours = int(remaining_seconds // 3600)
    minutes = int((remaining_seconds % 3600) // 60)
    seconds = int(remaining_seconds % 60)
    # Formatea el tiempo en formato HH:MM:SS como una cadena de texto
    return "{{:02d}}:{{:02d}}:{{:02d}}".format(hours, minutes, seconds)

# La actividad de Discord con la configuración del usuario
def discord_rpc_open(idClientDiscord, imageLarge, imageSmall, description_text, mainTexts_buttons):
    RPC = Presence(idClientDiscord)
    RPC.connect()

    start_time = time.time()

    # Bucle infinito para actualizar continuamente la presencia
    while True:
        RPC.update(
            state="Activity elapsed: " + get_current_time(start_time),
            large_image=imageLarge,
            small_image=imageSmall,
            details="Activity Summary:   " + description_text,
            party_size=[1, 1],
            buttons=mainTexts_buttons
        )

time.sleep(1)

idClientDiscord = "{idClientDiscord}"
imageLarge = "{image_large}"
imageSmall = "{image_small}"
description_text = "{description}"
mainTexts_buttons = {mainTexts_buttons}

discord_rpc_open(idClientDiscord, imageLarge, imageSmall, description_text, mainTexts_buttons)
    '''.format(idClientDiscord= idClientDiscord, 
            image_large=imageLarge, 
            image_small=imageSmall, 
            description=description_text, 
            mainTexts_buttons=mainTexts_buttons)

        save_file = int(input("Quieres guardar la configuracion, 1[Si] Otro[No]"))

        if save_file == 1:
            with open('file_discordRPC.pyw', 'w') as file:
                file.write(configuration)
                
            print("Configuracion Guardad con Exito. Hemos creado un archivo con esa configuracion.")
        

    else:
        invalid_option_text = Colorate.Horizontal(Colors.cyan_to_green, "              [!] Por favor elija una opcion valida.")
        print(invalid_option_text)
        time.sleep(2)

    time.sleep(.5)
    Banner_end()
#-----------------------------------------------------------------------------------------------------
    

#main 
if (__name__ == "__main__"):
    Banner()
    time.sleep(1)
    responses_functions()
